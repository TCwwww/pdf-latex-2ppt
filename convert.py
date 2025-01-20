import os
import sys
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches

def pdf_to_ppt(pdf_path, ppt_path, dpi=300):
    """
    Converts a PDF (created from LaTeX) into a PowerPoint presentation with each page as an image.

    Args:
        pdf_path (str): Path to the input PDF file.
        ppt_path (str): Path to save the output PPT file.
        dpi (int): DPI for converting PDF pages to images.
    """
    # Create a PowerPoint presentation
    presentation = Presentation()

    # Convert PDF pages to images
    images = convert_from_path(pdf_path, dpi=dpi)
    total_pages = len(images)
    
    print(f"Converting {total_pages} pages...")
    
    for idx, img in enumerate(images):
        # Show progress on same line
        progress = (idx + 1) / total_pages * 100
        print(f"Processing page {idx + 1}/{total_pages} ({progress:.1f}%)    ", end="\r")
        # Get image dimensions
        width, height = img.size
        
        if height > width:  # Vertical orientation
            # Split image into upper and lower halves
            upper_half = img.crop((0, 0, width, height//2))
            lower_half = img.crop((0, height//2, width, height))
            
            # Calculate aspect ratio for halves
            aspect_ratio = width / (height / 2)
            half_height = Inches(10) / aspect_ratio
            
            # Calculate vertical offset to center the image
            vertical_offset = (Inches(7.5) - half_height) / 2
            
            # Save and add upper half
            upper_path = f"temp_page_{idx + 1}_upper.png"
            upper_half.save(upper_path, "PNG")
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            slide.shapes.add_picture(upper_path, Inches(0), vertical_offset, width=Inches(10))
            os.remove(upper_path)
            
            # Save and add lower half
            lower_path = f"temp_page_{idx + 1}_lower.png"
            lower_half.save(lower_path, "PNG")
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            slide.shapes.add_picture(lower_path, Inches(0), vertical_offset, width=Inches(10))
            os.remove(lower_path)
        else:  # Horizontal or square orientation
            # Save and add full image
            full_path = f"temp_page_{idx + 1}.png"
            img.save(full_path, "PNG")
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            slide.shapes.add_picture(full_path, Inches(0), Inches(0), width=Inches(10))
            os.remove(full_path)

    # Save the presentation
    presentation.save(ppt_path)
    print(f"PowerPoint saved at {ppt_path}")

if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python convert.py input.pdf output.pptx")
        sys.exit(1)
        
    input_pdf = sys.argv[1]
    output_ppt = sys.argv[2]
    
    if not input_pdf.lower().endswith('.pdf'):
        print("Error: Input file must be a PDF")
        sys.exit(1)
        
    if not output_ppt.lower().endswith('.pptx'):
        print("Error: Output file must be a PPTX")
        sys.exit(1)
        
    pdf_to_ppt(input_pdf, output_ppt)
